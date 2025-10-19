from pptx import Presentation
import sys

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = f"\n{'='*80}\n{pptx_path}\n{'='*80}\n\n"
    
    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text += shape.text + "\n"
    
    return text

# Extract from all PPTX files
pptx_files = [
    "5-1-RASA Chatbot Frameworks.pptx",
    "6-1-Remote Interpreter Modul.pptx",
    "7-1-Integrasi dengan External API.pptx"
]

all_text = ""
for pptx_file in pptx_files:
    try:
        all_text += extract_text_from_pptx(pptx_file)
    except Exception as e:
        all_text += f"\nError reading {pptx_file}: {e}\n"

print(all_text)

# Save to file
with open("extracted_pptx_content.txt", "w", encoding="utf-8") as f:
    f.write(all_text)

print("\n\n*** Content saved to extracted_pptx_content.txt ***")
